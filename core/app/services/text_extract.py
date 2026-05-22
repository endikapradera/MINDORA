from __future__ import annotations

import io
import logging
import os
import re
import csv
from collections import Counter
from pathlib import Path

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader
from PIL import Image
import pytesseract

logger = logging.getLogger(__name__)

# Minimum characters per page to consider it "has text" and skip OCR
_OCR_THRESHOLD = 60
# Minimum image size (px) to attempt OCR / description
_MIN_IMG_PX = 80
_OCR_TIMEOUT_SECONDS = int(os.getenv("MINDORA_OCR_TIMEOUT_SECONDS", "8"))
_OCR_EMBEDDED_IMAGES = os.getenv("MINDORA_OCR_EMBEDDED_IMAGES", "0") == "1"
# A line present in more than this fraction of pages is considered a header/footer
_HEADER_FOOTER_THRESHOLD = 0.40


def extract_text_from_file(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix == ".docx":
        return _extract_docx(path)
    if suffix == ".pptx":
        return _extract_pptx(path)
    if suffix in {".png", ".jpg", ".jpeg", ".bmp", ".tiff"}:
        return _extract_image(path)
    if suffix == ".csv":
        return _extract_csv(path)
    if suffix in {".txt", ".md"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    raise ValueError("Unsupported file type")


# ---------------------------------------------------------------------------
# Header / footer deduplication
# ---------------------------------------------------------------------------

def _remove_repeated_lines(per_page_texts: list[str], threshold: float = _HEADER_FOOTER_THRESHOLD) -> list[str]:
    """
    Detect and strip lines that appear verbatim in more than `threshold` fraction
    of pages — these are almost certainly running headers, footers, or watermarks.

    Strategy:
      • Examine only the first 3 and last 3 lines of each page (where headers/footers live).
      • Count unique occurrences across pages (one count per page regardless of repeats within).
      • Remove any such line from ALL positions within each page's text.
    """
    n = len(per_page_texts)
    if n < 3:
        # Too few pages — not enough signal, skip dedup
        return per_page_texts

    # Collect candidate lines (first/last 3 of each page)
    candidate_counts: Counter[str] = Counter()
    for text in per_page_texts:
        lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
        seen_this_page: set[str] = set()
        candidates = lines[:3] + lines[-3:]
        for ln in candidates:
            if len(ln) >= 4 and ln not in seen_this_page:
                candidate_counts[ln] += 1
                seen_this_page.add(ln)

    # Anything exceeding the threshold is a header/footer
    repeated: set[str] = {
        ln for ln, count in candidate_counts.items()
        if count / n >= threshold
    }

    if not repeated:
        return per_page_texts

    logger.info(
        "Header/footer dedup: removing %d repeated line(s): %s",
        len(repeated),
        list(repeated)[:5],
    )

    cleaned: list[str] = []
    for text in per_page_texts:
        lines = text.splitlines()
        filtered = [ln for ln in lines if ln.strip() not in repeated]
        cleaned.append("\n".join(filtered))

    return cleaned


# ---------------------------------------------------------------------------
# PDF helpers
# ---------------------------------------------------------------------------

def _ocr_pil_image(img: Image.Image) -> str:
    """Run Tesseract OCR on a PIL image, return cleaned text."""
    try:
        text = pytesseract.image_to_string(
            img,
            lang="spa+eng",
            config="--psm 6",
            timeout=_OCR_TIMEOUT_SECONDS,
        )
        return text.strip()
    except BaseException as exc:  # pragma: no cover
        logger.warning("OCR error: %s", exc)
        return ""


def _fitz_page_to_pil(fitz_page) -> Image.Image:  # type: ignore[return]
    """Render a PyMuPDF page to a PIL image at 2× resolution for better OCR."""
    mat = __import__("fitz").Matrix(2.0, 2.0)
    pix = fitz_page.get_pixmap(matrix=mat, alpha=False)
    return Image.frombytes("RGB", [pix.width, pix.height], pix.samples)


def _describe_image_from_bytes(data: bytes, label: str) -> str:
    """
    Try to extract text from an embedded image via OCR.
    Returns a descriptive block like [IMAGEN: <ocr_text>].
    If OCR yields nothing useful, returns a generic placeholder.
    """
    try:
        img = Image.open(io.BytesIO(data)).convert("RGB")
        w, h = img.size
        if w < _MIN_IMG_PX or h < _MIN_IMG_PX:
            return ""
        ocr_text = _ocr_pil_image(img)
        if len(ocr_text) > 20:
            return f"[IMAGEN {label}: {ocr_text}]"
        # No readable text → mark as visual element (diagram/chart)
        aspect = w / h if h else 1.0
        kind = "diagrama" if 0.5 < aspect < 2.5 else "figura"
        return f"[IMAGEN {label}: {kind} sin texto extraíble, tamaño {w}×{h}px]"
    except Exception as exc:  # pragma: no cover
        logger.debug("Image description failed for %s: %s", label, exc)
        return ""


def _extract_pdf(path: Path) -> str:
    """
    Extract text from a PDF.
    - Pages with enough native text → use pypdf (fast).
    - Pages with little/no text → render with PyMuPDF and OCR with Tesseract.
    - Embedded images → OCR + describe inline as [IMAGEN n].
    """
    reader = PdfReader(str(path))

    # Try to import fitz (PyMuPDF) for OCR fallback; gracefully degrade if missing.
    try:
        import fitz  # type: ignore
        fitz_doc = fitz.open(str(path))
        _fitz_available = True
    except ImportError:  # pragma: no cover
        fitz_doc = None
        _fitz_available = False

    parts: list[str] = []
    img_counter = 0

    for i, page in enumerate(reader.pages, start=1):
        try:
            native_text = (page.extract_text() or "").strip()
        except BaseException as exc:  # pragma: no cover
            logger.warning("Native PDF extraction failed for page %d: %s", i, exc)
            native_text = ""

        # --- Native text path ---
        if len(native_text) >= _OCR_THRESHOLD:
            page_parts = [f"[PAGE {i}]\n{native_text}"]
        else:
            # --- OCR fallback path ---
            ocr_text = ""
            if _fitz_available and fitz_doc is not None:
                try:
                    fitz_page = fitz_doc[i - 1]
                    pil_img = _fitz_page_to_pil(fitz_page)
                    ocr_text = _ocr_pil_image(pil_img)
                    logger.info("PAGE %d: OCR aplicado (%d chars)", i, len(ocr_text))
                except BaseException as exc:  # pragma: no cover
                    logger.warning("OCR fallback failed for page %d: %s", i, exc)

            combined = (native_text + "\n" + ocr_text).strip()
            page_parts = [f"[PAGE {i}]\n{combined}"] if combined else []

        # --- Embedded images in page ---
        if _OCR_EMBEDDED_IMAGES:
            try:
                for img_obj in page.images:
                    img_counter += 1
                    raw = img_obj.data
                    desc = _describe_image_from_bytes(raw, f"p{i}-{img_counter}")
                    if desc:
                        page_parts.append(desc)
            except Exception:  # pragma: no cover
                pass  # pypdf may not expose images on all versions

        parts.extend(page_parts)

    if _fitz_available and fitz_doc is not None:
        fitz_doc.close()

    # Separate page markers from page bodies, dedup headers/footers, then rejoin
    # parts is a flat list of "[PAGE n]\n<text>" strings plus "[IMAGEN ...]" items
    page_blocks: list[str] = []
    image_stubs: list[tuple[int, str]] = []  # (original index, text)

    for idx, item in enumerate(parts):
        if item.startswith("[PAGE "):
            page_blocks.append(item)
        else:
            image_stubs.append((idx, item))

    # Strip headers/footers only from page body text (not from image descriptions)
    cleaned_pages = _remove_repeated_lines(page_blocks)

    # Rebuild parts in original order
    result_parts: list[str] = list(cleaned_pages)
    for orig_idx, stub in image_stubs:
        result_parts.append(stub)

    return "\n".join(result_parts)


def _extract_docx(path: Path) -> str:
    doc = DocxDocument(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text)


def _extract_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_parts.append(text)
        if slide_parts:
            parts.append(f"[SLIDE {i}]\n" + "\n".join(slide_parts))
    return "\n".join(parts)


def _extract_image(path: Path) -> str:
    try:
        with Image.open(path) as img:
            return pytesseract.image_to_string(img, lang="spa+eng")
    except pytesseract.TesseractNotFoundError as exc:
        raise ValueError(
            "Tesseract OCR no está instalado en el sistema. Instálalo para procesar imágenes."
        ) from exc


def _extract_csv(path: Path) -> str:
    """Extract CSV as structured text with schema, sample rows and numeric stats."""

    raw = ""
    for enc in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            raw = path.read_text(encoding=enc, errors="strict")
            break
        except Exception:
            continue
    if not raw:
        raw = path.read_text(encoding="utf-8", errors="ignore")

    reader = csv.DictReader(io.StringIO(raw))
    if not reader.fieldnames:
        return "[CSV] Archivo vacío o sin cabeceras válidas."

    headers = [h.strip() for h in reader.fieldnames if h and h.strip()]
    rows = [r for r in reader]
    total_rows = len(rows)

    sample_lines: list[str] = []
    for i, row in enumerate(rows[:8], start=1):
        compact = ", ".join(f"{k}={str(row.get(k, '')).strip()}" for k in headers[:8])
        sample_lines.append(f"Fila {i}: {compact}")

    numeric_stats: list[str] = []
    for col in headers:
        values: list[float] = []
        for r in rows:
            token = str(r.get(col, "")).strip().replace(",", ".")
            if not token:
                continue
            try:
                values.append(float(token))
            except Exception:
                continue
        if len(values) >= 2:
            avg = sum(values) / len(values)
            numeric_stats.append(
                f"- {col}: n={len(values)}, min={min(values):.2f}, max={max(values):.2f}, media={avg:.2f}"
            )

    stats_block = "\n".join(numeric_stats) if numeric_stats else "- No se detectaron columnas numéricas suficientes."
    sample_block = "\n".join(sample_lines) if sample_lines else "(sin filas de muestra)"

    return (
        "[CSV]\n"
        f"Columnas ({len(headers)}): {', '.join(headers)}\n"
        f"Total de filas: {total_rows}\n\n"
        "Estadísticas numéricas:\n"
        f"{stats_block}\n\n"
        "Muestra de datos:\n"
        f"{sample_block}"
    )
